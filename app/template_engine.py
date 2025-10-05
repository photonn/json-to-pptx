"""Utility helpers to render PowerPoint presentations based on JSON instructions.

The implementation borrows concepts from the original `pptx-template` project
but is intentionally minimal.  It replaces text placeholders wrapped in curly
braces and populates tables by name using row/column data provided in the
payload.  The module is designed to work with modern versions of Python and
`python-pptx` without pulling in the older runtime dependencies required by the
upstream project.
"""

from __future__ import annotations

import base64
import copy
import io
import re
from pathlib import Path
from dataclasses import dataclass, field
from typing import Any, Dict, Iterable, List, Mapping, Optional

from pptx import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.table import Table

PLACEHOLDER_PATTERN = re.compile(r"\{([A-Za-z0-9._\-]+)\}")
SLIDE_ID_PATTERN = re.compile(r"\{id:([A-Za-z0-9._\-]+)\}")


def _resolve_path(data: Any, path: str) -> Optional[Any]:
    """Resolve a dotted path inside ``data``."""

    parts = [part for part in path.split('.') if part]
    value: Any = data

    for part in parts:
        if isinstance(value, Mapping):
            value = value.get(part)
        elif isinstance(value, list) and part.isdigit():
            index = int(part)
            if 0 <= index < len(value):
                value = value[index]
            else:
                return None
        else:
            return None

        if value is None:
            return None

    return value


def _iter_text_shapes(slide) -> Iterable[Any]:
    for shape in slide.shapes:
        if not hasattr(shape, 'has_text_frame'):
            continue
        if shape.has_text_frame:
            yield shape


def _iter_tables(slide) -> Iterable[Table]:
    for shape in slide.shapes:
        if isinstance(shape, GraphicFrame) and shape.has_table:
            yield shape.table


def _replace_placeholders_in_text_frame(text_frame, model: Mapping[str, Any]):
    """Replace all ``{placeholder}`` entries with values from ``model``."""

    if not text_frame.text:
        return

    matches = list(PLACEHOLDER_PATTERN.finditer(text_frame.text))
    if not matches:
        return

    original_text = text_frame.text
    replaced_text = original_text

    for match in matches:
        path = match.group(1)
        value = _resolve_path(model, path)
        replacement = '' if value is None else str(value)
        replaced_text = replaced_text.replace(match.group(0), replacement)

    if replaced_text != original_text:
        text_frame.clear()
        text_frame.text = replaced_text


@dataclass
class TableInstruction:
    """Configuration for populating a table."""

    shape: str
    data: List[List[Any]]
    header: Optional[List[Any]] = None
    clear_extra_rows: bool = True


@dataclass
class SlideInstruction:
    """Instructions that apply to a single slide."""

    id: Optional[str] = None
    index: Optional[int] = None
    replacements: Mapping[str, Any] = field(default_factory=dict)
    tables: List[TableInstruction] = field(default_factory=list)

    def context(self, global_context: Mapping[str, Any]) -> Dict[str, Any]:
        merged = copy.deepcopy(global_context)
        merged.update(self.replacements)
        return merged


class TemplateEngine:
    """Render PowerPoint presentations using JSON instructions."""

    def __init__(self, templates_dir: str):
        self._templates_dir = Path(templates_dir)

    def render(self, payload: Mapping[str, Any]) -> bytes:
        template_name = payload.get('template', 'default')
        template_path = self._templates_dir / f"{template_name}.pptx"
        if not template_path.exists():
            raise ValueError(
                "Template '{0}' was not found in '{1}'".format(
                    template_name, self._templates_dir
                )
            )

        presentation = Presentation(str(template_path))
        context = payload.get('context', {})

        for slide_payload in payload.get('slides', []):
            instruction = self._parse_slide_instruction(slide_payload)
            slide = self._resolve_slide(presentation, instruction)
            model = instruction.context(context)

            for shape in _iter_text_shapes(slide):
                _replace_placeholders_in_text_frame(shape.text_frame, model)

            for table in _iter_tables(slide):
                for row in table.rows:
                    for cell in row.cells:
                        _replace_placeholders_in_text_frame(cell.text_frame, model)

            for table_instruction in instruction.tables:
                self._populate_table(slide, table_instruction, model)

        output = io.BytesIO()
        presentation.save(output)
        return output.getvalue()

    def _parse_slide_instruction(self, payload: Mapping[str, Any]) -> SlideInstruction:
        tables_payload = payload.get('tables', [])
        tables = [
            TableInstruction(
                shape=table_payload['shape'],
                data=table_payload.get('data', []),
                header=table_payload.get('header'),
                clear_extra_rows=table_payload.get('clear_extra_rows', True),
            )
            for table_payload in tables_payload
        ]

        return SlideInstruction(
            id=payload.get('id'),
            index=payload.get('index'),
            replacements=payload.get('replacements', {}),
            tables=tables,
        )

    def _resolve_slide(self, presentation: Presentation, instruction: SlideInstruction):
        if instruction.id:
            for slide in presentation.slides:
                for shape in _iter_text_shapes(slide):
                    text = shape.text
                    if not text:
                        continue
                    match = SLIDE_ID_PATTERN.search(text)
                    if match and match.group(1) == instruction.id:
                        return slide
            raise ValueError(f"Slide with id '{instruction.id}' was not found")

        if instruction.index is None:
            raise ValueError("Either 'id' or 'index' must be provided for a slide")

        try:
            return presentation.slides[instruction.index]
        except IndexError as exc:
            raise ValueError(
                f"Slide index {instruction.index} is out of range"
            ) from exc

    def _populate_table(
        self,
        slide,
        instruction: TableInstruction,
        model: Mapping[str, Any],
    ) -> None:
        target_shape = None
        for shape in slide.shapes:
            if getattr(shape, 'name', None) == instruction.shape and isinstance(
                shape, GraphicFrame
            ) and shape.has_table:
                target_shape = shape
                break

        if target_shape is None:
            raise ValueError(f"Table '{instruction.shape}' was not found on the slide")

        table = target_shape.table
        data_rows = copy.deepcopy(instruction.data)
        if instruction.header:
            data_rows.insert(0, instruction.header)

        required_rows = len(data_rows)
        current_rows = len(table.rows)

        if required_rows > current_rows:
            raise ValueError(
                f"Table '{instruction.shape}' requires at least {required_rows} rows, but only {current_rows} are present in the template."
            )

        if data_rows:
            required_cols = max(len(row) for row in data_rows)
        else:
            required_cols = len(table.columns)

        if required_cols > len(table.columns):
            raise ValueError(
                f"Table '{instruction.shape}' expects {len(table.columns)} columns"
            )

        for row_index, row_values in enumerate(data_rows):
            for col_index in range(len(table.columns)):
                cell = table.cell(row_index, col_index)
                value = row_values[col_index] if col_index < len(row_values) else ''
                cell.text = '' if value is None else str(value)

        if instruction.clear_extra_rows and required_rows < len(table.rows):
            for row_index in range(required_rows, len(table.rows)):
                for col_index in range(len(table.columns)):
                    cell = table.cell(row_index, col_index)
                    cell.text = ''


def encode_pptx(data: bytes) -> str:
    """Return a base64 representation of a PPTX file."""

    return base64.b64encode(data).decode('ascii')
